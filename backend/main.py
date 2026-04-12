"""
Wrytica document processor backend.
"""

from __future__ import annotations

import asyncio
import json
import logging
import os
import platform
import shutil
import subprocess
import tempfile
import time
import uuid
from importlib.metadata import PackageNotFoundError, version as package_version
from pathlib import Path
from typing import Any, Dict, List, Literal, Optional, Tuple

from fastapi import FastAPI, File, HTTPException, Query, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from pydantic import BaseModel

# Import resource and job management
from resource_manager import resource_monitor, ResourceMonitor, HardwareProfile, RESOURCE_LIMITS
from job_queue import job_queue, init_job_queue, JobStatus

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

APP_VERSION = "1.2.0"  # Updated for stability features
TEMP_DIR = Path(tempfile.gettempdir()) / "wrytica_backend"
TEMP_DIR.mkdir(exist_ok=True)

app = FastAPI(
    title="Wrytica Document Processor",
    description="Local backend for PDF, Office, and deep extraction processing.",
    version=APP_VERSION,
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=[],
    allow_origin_regex=r"^https?://(localhost|127\.0\.0\.1)(:\d+)?$",
    allow_credentials=False,
    allow_methods=["*"],
    allow_headers=["*"],
)

embedding_model = None
model_loading = False
_mineru_availability_cache: Optional[Dict[str, Any]] = None  # cached at startup


class DocumentChunk(BaseModel):
    id: str
    text: str
    page_number: Optional[int] = None
    section: Optional[str] = None
    word_count: int


class ProcessedDocument(BaseModel):
    document_id: str
    filename: str
    total_pages: Optional[int] = None
    total_chunks: int
    chunks: List[DocumentChunk]
    embeddings: Optional[List[List[float]]] = None
    processing_time_ms: float
    file_size_bytes: int


class DeepExtractResult(BaseModel):
    document_id: str
    filename: str
    markdown: str
    total_pages: Optional[int] = None
    processing_mode: Literal["gpu", "cpu", "mps", "fallback_pdfplumber"]
    processing_time_ms: float
    file_size_bytes: int
    layout_elements: Dict[str, int]


class HealthStatus(BaseModel):
    status: str
    version: str
    features: Dict[str, Any]


@app.on_event("startup")
async def startup_event():
    logger.info("Starting Wrytica backend")
    asyncio.create_task(load_embedding_model())


async def load_embedding_model():
    global embedding_model, model_loading
    if model_loading or embedding_model is not None:
        return

    model_loading = True
    try:
        from sentence_transformers import SentenceTransformer

        embedding_model = SentenceTransformer("all-MiniLM-L6-v2")
    except Exception as exc:
        logger.warning("Embedding model unavailable: %s", exc)
    finally:
        model_loading = False


def detect_compute_mode() -> Dict[str, Any]:
    info: Dict[str, Any] = {
        "platform": platform.system(),
        "machine": platform.machine().lower(),
        "accelerator": "cpu",
        "use_gpu": False,
        "free_vram_gb": 0.0,
        "reason": "cpu_default",
    }

    if info["platform"] == "Darwin" and info["machine"] in {"arm64", "aarch64"}:
        try:
            result = subprocess.run(
                ["sysctl", "-n", "hw.memsize"],
                capture_output=True,
                text=True,
                timeout=5,
                check=True,
            )
            total_memory_gb = int(result.stdout.strip()) / (1024**3)
            info["use_gpu"] = total_memory_gb >= 16.0
            info["accelerator"] = "mps" if info["use_gpu"] else "cpu"
            info["reason"] = "mps_available_unified_memory" if info["use_gpu"] else "insufficient_unified_memory"
        except Exception:
            info["reason"] = "mps_detection_failed"
        return info

    try:
        import pynvml

        pynvml.nvmlInit()
        handle = pynvml.nvmlDeviceGetHandleByIndex(0)
        mem = pynvml.nvmlDeviceGetMemoryInfo(handle)
        free_ratio = (mem.free / mem.total) if mem.total else 0.0
        info["free_vram_gb"] = round(mem.free / (1024**3), 2)
        info["use_gpu"] = info["free_vram_gb"] >= 4.0 and free_ratio >= 0.5
        info["accelerator"] = "gpu" if info["use_gpu"] else "cpu"
        info["reason"] = "gpu_available" if info["use_gpu"] else "gpu_busy_or_low_vram"
    except Exception:
        info["reason"] = "no_nvidia_vram_info"

    return info


def is_mineru_available() -> Dict[str, Any]:
    """Check MinerU availability. Result is cached after first call to avoid log spam."""
    global _mineru_availability_cache
    if _mineru_availability_cache is not None:
        return _mineru_availability_cache

    mineru_version: Optional[str] = None
    pipeline_ok = False
    try:
        mineru_version = package_version("mineru")
    except (PackageNotFoundError, Exception):
        mineru_version = None

    if mineru_version:
        try:
            from mineru.backend.pipeline.pipeline_analyze import doc_analyze_streaming  # noqa: F401
            from mineru.backend.pipeline.pipeline_middle_json_mkcontent import union_make  # noqa: F401
            pipeline_ok = True
            logger.info("MinerU %s pipeline ready", mineru_version)
        except Exception as e:
            # Log once at startup, then cache so it never repeats
            logger.warning("MinerU %s found but pipeline unavailable (dependency conflict): %s", mineru_version, str(e).split('\n')[0])
            logger.info("Advanced OCR will use Chandra (pypdfium2) as fallback")

    _mineru_availability_cache = {
        "available": pipeline_ok,
        "version": mineru_version,
        "pipeline_ok": pipeline_ok,
    }
    return _mineru_availability_cache


def _parse_content_list_metrics(items: List[Dict[str, Any]]) -> Dict[str, int]:
    text_like_types = {"text", "list", "code", "header", "footer", "page_number", "aside_text", "page_footnote"}
    max_page = max((int(item.get("page_idx", -1)) for item in items), default=-1)
    return {
        "total_pages": max_page + 1 if max_page >= 0 else 0,
        "text_blocks": sum(1 for item in items if item.get("type") in text_like_types),
        "tables": sum(1 for item in items if item.get("type") == "table"),
        "formulas": sum(1 for item in items if item.get("type") == "equation"),
        "images": sum(1 for item in items if item.get("type") == "image"),
        "figures": sum(1 for item in items if item.get("type") in {"image", "chart"}),
    }


def _fallback_markdown_from_content_list(items: List[Dict[str, Any]]) -> str:
    parts: List[str] = []
    for item in items:
        item_type = item.get("type")
        if item_type == "text":
            text = (item.get("text") or "").strip()
            if text:
                level = int(item.get("text_level") or 0)
                parts.append(f"{'#' * min(level + 1, 6)} {text}" if level > 0 else text)
        elif item_type == "equation":
            text = (item.get("text") or "").strip()
            if text:
                parts.append(text)
        elif item_type == "table":
            caption = " ".join(item.get("table_caption") or []).strip()
            body = (item.get("table_body") or item.get("text") or "").strip()
            if caption:
                parts.append(f"**Table:** {caption}")
            if body:
                parts.append(body)
    return "\n\n".join(parts)


def _load_mineru_output(output_dir: Path) -> Tuple[str, Dict[str, int]]:
    markdown_files = sorted(output_dir.rglob("*.md"), key=lambda p: p.stat().st_size if p.exists() else 0, reverse=True)
    markdown = markdown_files[0].read_text(encoding="utf-8", errors="ignore") if markdown_files else ""

    content_list_files = list(output_dir.rglob("*_content_list.json"))
    if content_list_files:
        items = json.loads(content_list_files[0].read_text(encoding="utf-8"))
        layout = _parse_content_list_metrics(items)
        if not markdown.strip():
            markdown = _fallback_markdown_from_content_list(items)
        return markdown, layout

    return markdown, {
        "total_pages": 0,
        "text_blocks": 0,
        "tables": 0,
        "formulas": 0,
        "images": 0,
        "figures": 0,
    }


def _run_mineru_pipeline(file_path: Path, compute: Dict[str, Any]) -> Tuple[str, Dict[str, int], str]:
    """Use MinerU Python API directly (bypasses CLI and office dependencies)."""
    from mineru.backend.pipeline.pipeline_analyze import doc_analyze_streaming
    from mineru.backend.pipeline.pipeline_middle_json_mkcontent import union_make
    from mineru.data.data_reader_writer import FileBasedDataWriter

    output_dir = TEMP_DIR / f"mineru_{uuid.uuid4().hex}"
    output_dir.mkdir(parents=True, exist_ok=True)
    image_dir = output_dir / "images"
    image_dir.mkdir(exist_ok=True)

    # Keep the reported mode aligned with the API schema and frontend labels.
    mode = compute.get("accelerator", "cpu")

    try:
        pdf_bytes = file_path.read_bytes()
        image_writer = FileBasedDataWriter(str(image_dir))
        middle_json_result: List[Any] = []

        def on_doc_ready(pdf_info_dict: Any) -> None:
            middle_json_result.append(pdf_info_dict)

        doc_analyze_streaming(
            pdf_bytes_list=[pdf_bytes],
            image_writer_list=[image_writer],
            lang_list=["en"],
            on_doc_ready=on_doc_ready,
            parse_method="auto",
            formula_enable=True,
            table_enable=True,
        )

        if not middle_json_result:
            raise RuntimeError("MinerU produced no output")

        markdown = union_make(middle_json_result[0], "markdown", str(image_dir))
        if not markdown.strip():
            raise RuntimeError("MinerU produced empty markdown")

        layout = _parse_content_list_metrics(middle_json_result[0])
        return markdown, layout, mode

    finally:
        shutil.rmtree(output_dir, ignore_errors=True)


async def _extract_with_mineru(file_path: Path, compute: Dict[str, Any]) -> Tuple[str, Dict[str, int], str]:
    return await asyncio.to_thread(_run_mineru_pipeline, file_path, compute)


def _extract_pdf_pages_as_markdown(file_path: Path) -> Tuple[List[str], int, int]:
    try:
        import pdfplumber

        parts: List[str] = []
        text_blocks = 0
        with pdfplumber.open(file_path) as pdf:
            total_pages = len(pdf.pages)
            for page_num, page in enumerate(pdf.pages, 1):
                cleaned = (page.extract_text() or "").strip()
                if cleaned:
                    text_blocks += 1
                    parts.append(f"## Page {page_num}\n\n{cleaned}")
        return parts, total_pages, text_blocks
    except ImportError:
        logger.warning("pdfplumber is not installed, falling back to PyPDF2")

    try:
        from PyPDF2 import PdfReader
    except ImportError as exc:
        raise RuntimeError(
            "No PDF fallback parser is installed. Install backend requirements with: pip install -r backend/requirements.txt"
        ) from exc

    reader = PdfReader(str(file_path))
    parts = []
    text_blocks = 0
    total_pages = len(reader.pages)
    for page_num, page in enumerate(reader.pages, 1):
        cleaned = (page.extract_text() or "").strip()
        if cleaned:
            text_blocks += 1
            parts.append(f"## Page {page_num}\n\n{cleaned}")
    return parts, total_pages, text_blocks


def _extract_with_pdf_fallback_sync(file_path: Path) -> Tuple[str, Dict[str, int], str]:
    parts, total_pages, text_blocks = _extract_pdf_pages_as_markdown(file_path)
    return (
        "\n\n---\n\n".join(parts),
        {
            "total_pages": total_pages,
            "text_blocks": text_blocks,
            "tables": 0,
            "formulas": 0,
            "images": 0,
            "figures": 0,
        },
        "fallback_pdfplumber",
    )


async def _extract_with_pdf_fallback(file_path: Path) -> Tuple[str, Dict[str, int], str]:
    return await asyncio.to_thread(_extract_with_pdf_fallback_sync, file_path)


def _extract_with_chandra_sync(file_path: Path) -> Tuple[str, Dict[str, int], str]:
    """
    Chandra engine: balanced OCR using pypdfium2 for high-fidelity text extraction.
    Better layout awareness than pdfplumber, works without GPU.
    """
    try:
        import pypdfium2 as pdfium
    except ImportError:
        logger.warning("pypdfium2 not available for Chandra, falling back to pdfplumber")
        return _extract_with_pdf_fallback_sync(file_path)

    try:
        pdf = pdfium.PdfDocument(str(file_path))
        parts: List[str] = []
        text_blocks = 0
        tables = 0
        images_found = 0
        total_pages = len(pdf)

        for page_num in range(total_pages):
            page = pdf[page_num]
            page_label = f"Page {page_num + 1}"

            # Extract text with layout awareness via pypdfium2 text page
            textpage = page.get_textpage()

            # Get full page text with position-aware extraction
            page_text = textpage.get_text_range()

            # Count images on the page
            try:
                img_count = sum(1 for _ in page.get_objects(filter=[pdfium.raw.FPDF_PAGEOBJ_IMAGE]))
            except Exception:
                img_count = 0
            images_found += img_count

            # Detect likely tables by counting lines with tab/column separators
            lines = page_text.split('\n')
            tab_lines = sum(1 for ln in lines if '\t' in ln or '  ' in ln.strip())
            if tab_lines > 3:
                tables += 1
                # Format as code block to preserve alignment
                formatted = f"## {page_label}\n\n```\n{page_text.strip()}\n```"
            else:
                # Clean up excessive whitespace but preserve paragraph structure
                cleaned_lines = []
                for ln in lines:
                    stripped = ln.rstrip()
                    cleaned_lines.append(stripped)

                # Collapse multiple blank lines to single
                import re
                cleaned = re.sub(r'\n{3,}', '\n\n', '\n'.join(cleaned_lines)).strip()

                if cleaned:
                    formatted = f"## {page_label}\n\n{cleaned}"
                    text_blocks += 1
                else:
                    formatted = f"## {page_label}\n\n*(No text extracted — page may be image-based)*"

            if img_count > 0:
                formatted += f"\n\n*[{img_count} image(s) on this page]*"

            parts.append(formatted)
            textpage.close()
            page.close()

        pdf.close()

        return (
            "\n\n---\n\n".join(parts),
            {
                "total_pages": total_pages,
                "text_blocks": text_blocks,
                "tables": tables,
                "formulas": 0,
                "images": images_found,
                "figures": images_found,
            },
            "chandra_pypdfium2",
        )

    except Exception as e:
        logger.warning(f"Chandra engine failed: {e}, falling back to pdfplumber")
        return _extract_with_pdf_fallback_sync(file_path)


async def _extract_with_chandra(file_path: Path) -> Tuple[str, Dict[str, int], str]:
    return await asyncio.to_thread(_extract_with_chandra_sync, file_path)


def chunk_text(text: str, chunk_size: int = 800, overlap: int = 200) -> List[str]:
    if len(text) <= chunk_size:
        return [text] if text.strip() else []

    chunks: List[str] = []
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        if end < len(text):
            for idx in range(end - 1, max(start, end - 100), -1):
                if idx < len(text) and text[idx] in ".!?" and idx + 1 < len(text) and text[idx + 1] in " \n":
                    end = idx + 1
                    break
            else:
                while end > start and end < len(text) and text[end - 1] not in " \n":
                    end -= 1
                if end <= start:
                    end = min(start + chunk_size, len(text))
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end >= len(text):
            break
        next_start = max(0, end - overlap)
        start = next_start if next_start > start else end
    return chunks


def _build_deep_extract_chunks(markdown: str) -> List[Dict[str, Any]]:
    chunks: List[Dict[str, Any]] = []
    current_section: Optional[str] = None
    for idx, piece in enumerate(chunk_text(markdown, chunk_size=1200, overlap=150)):
        for line in piece.splitlines():
            stripped = line.strip()
            if stripped.startswith("#"):
                current_section = stripped.lstrip("#").strip() or current_section
                break
        chunks.append(
            {
                "id": f"chunk-deep-{idx}",
                "text": piece,
                "page_number": None,
                "section": current_section,
                "word_count": len(piece.split()),
            }
        )
    return chunks


async def deep_extract_pdf_file(file_path: Path, filename: str, file_size_bytes: int) -> DeepExtractResult:
    start_time = time.time()
    mineru_info = is_mineru_available()
    compute = detect_compute_mode()

    try:
        if mineru_info["pipeline_ok"]:
            markdown, layout, mode = await _extract_with_mineru(file_path, compute)
        else:
            logger.info("MinerU pipeline unavailable, falling back to pdfplumber")
            markdown, layout, mode = await _extract_with_pdf_fallback(file_path)
    except Exception as exc:
        logger.warning("MinerU extraction failed, falling back to pdfplumber: %s", exc)
        markdown, layout, mode = await _extract_with_pdf_fallback(file_path)

    return DeepExtractResult(
        document_id=str(uuid.uuid4()),
        filename=filename,
        markdown=markdown,
        total_pages=layout.get("total_pages"),
        processing_mode=mode,
        processing_time_ms=(time.time() - start_time) * 1000,
        file_size_bytes=file_size_bytes,
        layout_elements={
            "text_blocks": int(layout.get("text_blocks", 0)),
            "tables": int(layout.get("tables", 0)),
            "formulas": int(layout.get("formulas", 0)),
            "images": int(layout.get("images", 0)),
            "figures": int(layout.get("figures", 0)),
        },
    )


@app.on_event("startup")
async def startup_event():
    """Initialize backend services on startup."""
    logger.info("Starting Wrytica backend v%s", APP_VERSION)

    # Initialize job queue with resource monitor
    global job_queue
    profile = resource_monitor.detect_profile()
    resource_monitor.profile = profile
    resource_monitor.limits = RESOURCE_LIMITS[profile]

    max_concurrent = resource_monitor.limits["max_concurrent_jobs"]
    job_queue = init_job_queue(max_concurrent, resource_monitor)

    logger.info(f"Resource monitor initialized: profile={profile}, max_concurrent={max_concurrent}")

    # Cache MinerU availability once at startup (prevents repeated warnings on every health poll)
    is_mineru_available()

    # Start background job worker
    asyncio.create_task(load_embedding_model())
    asyncio.create_task(job_queue.process_worker(ocr_job_handler))


@app.get("/health", response_model=HealthStatus)
async def health_check():
    mineru_info = is_mineru_available()
    compute = detect_compute_mode()
    # Check which OCR engines are available
    chandra_available = False
    try:
        import pypdfium2  # noqa: F401
        chandra_available = True
    except ImportError:
        pass
    pdfplumber_available = False
    try:
        import pdfplumber  # noqa: F401
        pdfplumber_available = True
    except ImportError:
        pass
    ocr_available = chandra_available or pdfplumber_available or bool(mineru_info["available"])
    return HealthStatus(
        status="healthy",
        version=APP_VERSION,
        features={
            "pdf_processing": True,
            "office_processing": True,
            "embeddings": embedding_model is not None,
            "ocr": ocr_available,
            "ocr_fast": pdfplumber_available,
            "ocr_balanced": chandra_available,
            "deep_extract": bool(mineru_info["available"]),
            "deep_extract_gpu": bool(mineru_info["available"] and compute.get("accelerator") in {"gpu", "mps"}),
            "deep_extract_cpu": bool(mineru_info["available"]),
            "mineru_version": mineru_info.get("version"),
            "deep_extract_compute_reason": compute.get("reason"),
        },
    )


# ==============================================================================
# Job Queue & Resource Monitoring Endpoints
# ==============================================================================

@app.get("/api/system/metrics")
async def get_system_metrics():
    """Get current system resource usage."""
    snapshot = resource_monitor.get_snapshot()
    queue_stats = job_queue.get_queue_stats() if job_queue else {}

    return {
        "resources": snapshot.to_dict(),
        "queue": queue_stats,
        "hardware_profile": resource_monitor.profile.value,
    }


@app.get("/api/system/stats")
async def get_system_stats():
    """Get detailed system statistics and job history."""
    return resource_monitor.get_stats()


@app.post("/api/jobs/ocr")
async def start_ocr_job(file: UploadFile = File(...), engine: str = Query("auto"), timeout_sec: int = Query(1800)):
    """Start OCR job asynchronously with resource checking."""
    if not file.filename:
        raise HTTPException(status_code=400, detail="No filename provided")

    # Check file type
    filename_lower = file.filename.lower()
    if not any(filename_lower.endswith(ext) for ext in ['.pdf', '.png', '.jpg', '.jpeg', '.tiff', '.tif']):
        raise HTTPException(status_code=400, detail="Unsupported file type. Supported: PDF, PNG, JPG, TIFF")

    # Read file
    content = await file.read()
    if not content:
        raise HTTPException(status_code=400, detail="Uploaded file is empty")

    file_size_mb = len(content) / (1024 * 1024)

    # Check if system can handle it (non-blocking)
    job_id = str(uuid.uuid4())
    task_type = f"ocr_{engine}" if engine != "auto" else "ocr_auto"

    can_start, reason = resource_monitor.can_start_job(task_type)
    if not can_start:
        logger.warning(f"Job {job_id} cannot start: {reason}")

    # Save uploaded file
    temp_path = TEMP_DIR / f"{job_id}_{Path(file.filename).name}"
    try:
        temp_path.write_bytes(content)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to save file: {e}")

    # Queue job
    try:
        await job_queue.add_job(
            job_id,
            task_type,
            {"file_path": str(temp_path), "engine": engine},
            timeout_sec=timeout_sec,
            file_size_mb=file_size_mb
        )
    except Exception as e:
        logger.error(f"Failed to queue job: {e}")
        temp_path.unlink(missing_ok=True)
        raise HTTPException(status_code=500, detail=f"Failed to queue job: {e}")

    return {
        "job_id": job_id,
        "status": "queued",
        "file_size_mb": round(file_size_mb, 2),
        "can_start_immediately": can_start,
    }


@app.get("/api/jobs/{job_id}")
async def get_job_status(job_id: str):
    """Poll job status (non-blocking)."""
    job = job_queue.get_status(job_id)

    if not job:
        raise HTTPException(status_code=404, detail="Job not found")

    return job.to_dict()


@app.get("/api/jobs")
async def list_jobs(status: Optional[str] = Query(None)):
    """List all jobs, optionally filtered by status."""
    all_jobs = job_queue.get_all_jobs()

    jobs_list = [job.to_dict() for job in all_jobs.values()]

    if status:
        jobs_list = [j for j in jobs_list if j["status"] == status]

    return {"total": len(jobs_list), "jobs": jobs_list}


@app.post("/api/jobs/{job_id}/cancel")
async def cancel_job(job_id: str):
    """Cancel a job."""
    cancelled = await job_queue.cancel_job(job_id)

    if not cancelled:
        job = job_queue.get_status(job_id)
        if not job:
            raise HTTPException(status_code=404, detail="Job not found")

        return {
            "job_id": job_id,
            "cancelled": False,
            "reason": f"Cannot cancel job in status: {job.status.value}",
        }

    return {"job_id": job_id, "cancelled": True}


# ==============================================================================
# Background Job Handler
# ==============================================================================

async def ocr_job_handler(job_id: str, task_type: str, payload: Dict[str, Any]) -> Dict[str, Any]:
    """Handle OCR job execution."""
    file_path = Path(payload["file_path"])
    engine = payload.get("engine", "auto")

    logger.info(f"Processing job {job_id} with engine {engine}, file: {file_path.name}")

    start_time = time.time()

    try:
        # Update progress
        await job_queue.set_progress(job_id, 10, remaining_sec=None)

        # Determine which extraction engine to use
        if engine == "auto":
            mineru_available = is_mineru_available()["available"]
            compute = detect_compute_mode()

            # Auto-select based on hardware and file
            if mineru_available and compute["accelerator"] in {"gpu", "mps"}:
                engine = "mineru"
            else:
                engine = "pdfplumber"

        # Execute OCR based on engine
        markdown = ""
        layout_elements = {
            "text_blocks": 0,
            "tables": 0,
            "formulas": 0,
            "images": 0,
            "figures": 0,
        }
        mode = "fallback_pdfplumber"

        if engine == "mineru":
            mineru_info = is_mineru_available()
            if mineru_info["pipeline_ok"]:
                try:
                    await job_queue.set_progress(job_id, 20)
                    markdown, layout, mode = await _extract_with_mineru(file_path, detect_compute_mode())
                    layout_elements = layout
                except Exception as e:
                    logger.warning(f"MinerU pipeline failed for {job_id}, falling back to Chandra: {e}")
                    markdown, layout, mode = await _extract_with_chandra(file_path)
                    layout_elements = layout
            else:
                logger.info(f"MinerU not ready (dependency issue), using Chandra for job {job_id}")
                await job_queue.set_progress(job_id, 20)
                markdown, layout, mode = await _extract_with_chandra(file_path)
                layout_elements = layout
                mode = "chandra_pypdfium2 (mineru_unavailable)"

        elif engine == "chandra":
            try:
                await job_queue.set_progress(job_id, 20)
                markdown, layout, mode = await _extract_with_chandra(file_path)
                layout_elements = layout
            except Exception as e:
                logger.warning(f"Chandra failed for {job_id}, falling back to pdfplumber: {e}")
                markdown, layout, mode = await _extract_with_pdf_fallback(file_path)
                layout_elements = layout

        else:
            # pdfplumber (fast mode)
            await job_queue.set_progress(job_id, 20)
            markdown, layout, mode = await _extract_with_pdf_fallback(file_path)
            layout_elements = layout

        await job_queue.set_progress(job_id, 90)

        # Build chunks for knowledge base
        chunks = _build_deep_extract_chunks(markdown)

        await job_queue.set_progress(job_id, 100)

        duration_sec = time.time() - start_time

        logger.info(f"Job {job_id} completed in {duration_sec:.1f}s using {engine}")

        return {
            "document_id": str(uuid.uuid4()),
            "filename": file_path.name,
            "markdown": markdown,
            "chunks": chunks,
            "layout_elements": layout_elements,
            "processing_mode": mode,
            "processing_time_ms": duration_sec * 1000,
            "engine": engine,
        }

    except Exception as e:
        logger.error(f"Job {job_id} failed: {e}", exc_info=True)
        raise

    finally:
        # Cleanup temp file
        try:
            file_path.unlink(missing_ok=True)
        except Exception as e:
            logger.warning(f"Failed to clean up {file_path}: {e}")


@app.post("/api/v1/ocr/deep-extract", response_model=DeepExtractResult)
async def deep_extract_pdf(file: UploadFile = File(...)):
    if not file.filename or not file.filename.lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Only PDF files are supported.")

    content = await file.read()
    if not content:
        raise HTTPException(status_code=400, detail="Uploaded PDF is empty.")

    temp_file = TEMP_DIR / f"{uuid.uuid4()}_{Path(file.filename).name}"
    try:
        temp_file.write_bytes(content)
        return await deep_extract_pdf_file(temp_file, file.filename, len(content))
    except Exception as exc:
        logger.error("Deep extract failed for %s: %s", file.filename, exc)
        raise HTTPException(status_code=500, detail=f"Deep extract failed: {exc}") from exc
    finally:
        temp_file.unlink(missing_ok=True)


@app.post("/api/documents/process", response_model=ProcessedDocument)
async def process_document(
    file: UploadFile = File(...),
    include_embeddings: bool = Query(False),
    extraction_mode: Literal["standard", "deep"] = Query("standard"),
):
    start_time = time.time()
    logger.info("Processing file: %s (%s bytes), extraction_mode=%s", file.filename, file.size, extraction_mode)

    if not file.filename:
        raise HTTPException(status_code=400, detail="No filename provided")

    content = await file.read()
    temp_file = TEMP_DIR / f"{uuid.uuid4()}_{Path(file.filename).name}"
    temp_file.write_bytes(content)

    try:
        extension = Path(file.filename).suffix.lower()
        total_pages: Optional[int] = None

        if extension == ".pdf":
            chunks, total_pages = await process_pdf(temp_file, extraction_mode)
        elif extension in [".docx", ".doc"]:
            chunks = await process_docx(temp_file)
        elif extension in [".xlsx", ".xls"]:
            chunks = await process_excel(temp_file)
        elif extension in [".pptx", ".ppt"]:
            chunks = await process_pptx(temp_file)
        else:
            chunks = await process_text(temp_file)

        embeddings = None
        if include_embeddings and embedding_model and chunks:
            embeddings = embedding_model.encode([chunk["text"] for chunk in chunks]).tolist()

        return ProcessedDocument(
            document_id=str(uuid.uuid4()),
            filename=file.filename,
            total_pages=total_pages,
            total_chunks=len(chunks),
            chunks=[DocumentChunk(**chunk) for chunk in chunks],
            embeddings=embeddings,
            processing_time_ms=(time.time() - start_time) * 1000,
            file_size_bytes=len(content),
        )
    except Exception as exc:
        logger.error("Error processing %s: %s", file.filename, exc)
        raise HTTPException(status_code=500, detail=f"Processing error: {exc}") from exc
    finally:
        temp_file.unlink(missing_ok=True)


async def process_pdf(file_path: Path, extraction_mode: Literal["standard", "deep"]) -> Tuple[List[Dict[str, Any]], Optional[int]]:
    if extraction_mode == "deep":
        deep_result = await deep_extract_pdf_file(file_path, file_path.name, file_path.stat().st_size)
        return _build_deep_extract_chunks(deep_result.markdown), deep_result.total_pages

    parts, total_pages, _ = await asyncio.to_thread(_extract_pdf_pages_as_markdown, file_path)
    chunks: List[Dict[str, Any]] = []
    for page_number, part in enumerate(parts, 1):
        for idx, piece in enumerate(chunk_text(part, chunk_size=800, overlap=200)):
            chunks.append(
                {
                    "id": f"chunk-{page_number}-{idx}",
                    "text": piece,
                    "page_number": page_number,
                    "section": None,
                    "word_count": len(piece.split()),
                }
            )
    return chunks, total_pages


async def process_docx(file_path: Path) -> List[Dict[str, Any]]:
    from docx import Document

    document = Document(file_path)
    text = "\n".join(para.text for para in document.paragraphs if para.text.strip())
    return [
        {
            "id": f"chunk-{idx}",
            "text": chunk,
            "page_number": None,
            "section": None,
            "word_count": len(chunk.split()),
        }
        for idx, chunk in enumerate(chunk_text(text))
    ]


async def process_excel(file_path: Path) -> List[Dict[str, Any]]:
    import openpyxl

    workbook = openpyxl.load_workbook(file_path, data_only=True)
    blocks: List[str] = []
    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        rows = [f"Sheet: {sheet_name}"]
        for row in sheet.iter_rows():
            row_text = " | ".join(str(cell.value) for cell in row if cell.value)
            if row_text:
                rows.append(row_text)
        blocks.append("\n".join(rows))
    text = "\n\n".join(blocks)
    return [
        {
            "id": f"chunk-{idx}",
            "text": chunk,
            "page_number": None,
            "section": None,
            "word_count": len(chunk.split()),
        }
        for idx, chunk in enumerate(chunk_text(text))
    ]


async def process_pptx(file_path: Path) -> List[Dict[str, Any]]:
    from pptx import Presentation

    presentation = Presentation(file_path)
    blocks: List[str] = []
    for slide_num, slide in enumerate(presentation.slides, 1):
        slide_text = [f"Slide {slide_num}"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
        blocks.append("\n".join(slide_text))
    text = "\n\n".join(blocks)
    return [
        {
            "id": f"chunk-{idx}",
            "text": chunk,
            "page_number": None,
            "section": None,
            "word_count": len(chunk.split()),
        }
        for idx, chunk in enumerate(chunk_text(text))
    ]


async def process_text(file_path: Path) -> List[Dict[str, Any]]:
    text = file_path.read_text(encoding="utf-8", errors="ignore")
    return [
        {
            "id": f"chunk-{idx}",
            "text": chunk,
            "page_number": None,
            "section": None,
            "word_count": len(chunk.split()),
        }
        for idx, chunk in enumerate(chunk_text(text))
    ]


@app.post("/api/embeddings/generate")
async def generate_embeddings(texts: List[str]):
    if not embedding_model:
        raise HTTPException(status_code=503, detail="Embedding model not loaded yet")
    try:
        embeddings = embedding_model.encode(texts).tolist()
        return {"embeddings": embeddings, "dimensions": len(embeddings[0]) if embeddings else 0}
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"Embedding error: {exc}") from exc


if __name__ == "__main__":
    import uvicorn

    uvicorn.run(app, host="0.0.0.0", port=8000)
